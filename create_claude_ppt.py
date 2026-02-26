"""
Claude AI 생태계 완전 가이드 PPT 생성 스크립트
Claude's Brand Color Scheme + Clean Professional Design
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Claude Brand Colors ───
CLAUDE_CORAL = RGBColor(0xDA, 0x77, 0x56)      # #DA7756 - Claude 메인 코랄
CLAUDE_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # #1A1A2E - 다크 배경
CLAUDE_NAVY = RGBColor(0x2D, 0x2B, 0x55)       # #2D2B55 - 네이비
CLAUDE_CREAM = RGBColor(0xFA, 0xF3, 0xE8)      # #FAF3E8 - 크림
CLAUDE_LIGHT = RGBColor(0xF5, 0xF0, 0xEB)      # #F5F0EB - 밝은 배경
CLAUDE_BLUE = RGBColor(0x5B, 0x8D, 0xEF)       # #5B8DEF - 액센트 블루
CLAUDE_GREEN = RGBColor(0x4E, 0xC9, 0xB0)      # #4EC9B0 - 액센트 그린
CLAUDE_ORANGE = RGBColor(0xE8, 0x91, 0x4A)     # #E8914A - 액센트 오렌지
CLAUDE_PURPLE = RGBColor(0x9B, 0x72, 0xCF)     # #9B72CF - 액센트 퍼플
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x6B, 0x7B)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xEC)
DARK_GRAY = RGBColor(0x3A, 0x3A, 0x4A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Helper Functions ───

def add_dark_bg(slide):
    """다크 배경 추가"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = CLAUDE_DARK

def add_light_bg(slide):
    """밝은 배경 추가"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def add_cream_bg(slide):
    """크림 배경"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = CLAUDE_CREAM

def add_shape(slide, left, top, width, height, color, radius=None):
    """색상 박스 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_circle(slide, left, top, size, color):
    """원형 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=14,
                   default_color=BLACK, line_spacing=1.3):
    """여러 줄 텍스트 (리스트 형태)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text = line_data.get("text", "")
        p.text = text
        p.font.size = Pt(line_data.get("size", default_size))
        p.font.color.rgb = line_data.get("color", default_color)
        p.font.bold = line_data.get("bold", False)
        p.font.name = line_data.get("font", "맑은 고딕")
        p.alignment = line_data.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(line_data.get("space_after", 4))
        if line_data.get("indent"):
            p.level = line_data["indent"]
    return txBox

def add_accent_line(slide, left, top, width, color=CLAUDE_CORAL):
    """액센트 라인 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, items, icon_color=CLAUDE_CORAL,
             bg_color=WHITE, title_color=BLACK, text_color=GRAY):
    """카드 스타일 박스"""
    # 카드 배경
    card = add_shape(slide, left, top, width, height, bg_color)
    # 상단 액센트
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = icon_color
    accent.line.fill.background()
    # 타이틀
    add_text_box(slide, left + Inches(0.3), top + Inches(0.25), width - Inches(0.6), Inches(0.5),
                 title, font_size=16, color=title_color, bold=True)
    # 내용
    lines = [{"text": item, "size": 12, "color": text_color, "space_after": 3} for item in items]
    add_multi_text(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6),
                   height - Inches(0.9), lines, default_color=text_color)
    return card


# ══════════════════════════════════════════════════════════════
# SLIDE 1: 표지
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_dark_bg(slide)

# 왼쪽 큰 장식 원
add_circle(slide, Inches(-1), Inches(-1), Inches(4), CLAUDE_NAVY)
add_circle(slide, Inches(10), Inches(4.5), Inches(5), CLAUDE_NAVY)

# 중앙 코랄 액센트 라인
add_shape(slide, Inches(4.5), Inches(1.5), Inches(0.8), Pt(5), CLAUDE_CORAL)

# 타이틀
add_text_box(slide, Inches(2), Inches(2), Inches(9), Inches(1.2),
             "Claude AI 생태계 완전 가이드", font_size=42, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.3), Inches(9), Inches(0.8),
             "입문부터 고급 활용까지 — Code · Web · Desktop · Cowork · OpenClaw",
             font_size=18, color=CLAUDE_CORAL, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.3), Inches(9), Inches(0.5),
             "2026년 2월 최신판", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# 하단 정보
add_text_box(slide, Inches(2), Inches(6.2), Inches(9), Inches(0.5),
             "Made by Bella (OZKIZ) + Claude Opus 4.6  |  Anthropic Claude Max",
             font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: 목차
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5))
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(5), Inches(0.7),
             "목차", font_size=32, color=BLACK, bold=True)

toc_items = [
    ("01", "Claude 생태계 한눈에 보기", CLAUDE_CORAL),
    ("02", "Claude AI (Web)", CLAUDE_BLUE),
    ("03", "Claude Code (CLI)", CLAUDE_GREEN),
    ("04", "Claude Desktop & Cowork", CLAUDE_PURPLE),
    ("05", "Remote Control (NEW!)", CLAUDE_ORANGE),
    ("06", "제품별 비교표", CLAUDE_CORAL),
    ("07", "OpenClaw — 실행형 AI 에이전트", CLAUDE_BLUE),
    ("08", "OpenClaw 카카오톡 연동", CLAUDE_GREEN),
    ("09", "자동화 활용 — 입문자", CLAUDE_PURPLE),
    ("10", "자동화 활용 — 중급/고급", CLAUDE_ORANGE),
    ("11", "결합 활용 맵", CLAUDE_CORAL),
    ("12", "프롬프트 팁 & 베스트 프랙티스", CLAUDE_BLUE),
    ("13", "비용 & 시작 가이드", CLAUDE_GREEN),
]

for i, (num, title, color) in enumerate(toc_items):
    y = Inches(1.7) + Inches(i * 0.4)
    col = Inches(0.8) if i < 7 else Inches(6.8)
    y_adj = y if i < 7 else Inches(1.7) + Inches((i - 7) * 0.4)
    # 번호 원
    circle = add_circle(slide, col, y_adj, Inches(0.3), color)
    add_text_box(slide, col + Pt(2), y_adj + Pt(2), Inches(0.3), Inches(0.3),
                 num, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 텍스트
    add_text_box(slide, col + Inches(0.4), y_adj, Inches(5), Inches(0.35),
                 title, font_size=14, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: Claude 생태계 개요
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_shape(slide, Inches(0.8), Inches(0.5), Inches(1.2), Pt(4), CLAUDE_CORAL)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.7),
             "Claude 생태계 한눈에 보기", font_size=30, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.5),
             "Anthropic이 만든 AI 어시스턴트 Claude — 다양한 형태로 만나보세요",
             font_size=14, color=GRAY)

# 5개 제품 카드
products = [
    ("Claude AI\n(Web)", "claude.ai에서\n브라우저로 대화", CLAUDE_CORAL, "가장 기본"),
    ("Claude\nCode", "터미널 기반\nAI 코딩 에이전트", CLAUDE_GREEN, "개발자용"),
    ("Claude\nDesktop", "PC 앱에서\nAI와 작업", CLAUDE_BLUE, "데스크톱"),
    ("Claude\nCowork", "실제 업무 도구와\n직접 연동", CLAUDE_PURPLE, "업무 자동화"),
    ("Open\nClaw", "내 PC에서 24/7\nAI 비서 운영", CLAUDE_ORANGE, "자가호스팅"),
]

for i, (name, desc, color, badge) in enumerate(products):
    x = Inches(0.6) + Inches(i * 2.5)
    y = Inches(2.3)

    # 카드 배경
    card = add_shape(slide, x, y, Inches(2.2), Inches(3.8), CLAUDE_NAVY)

    # 상단 색상 바
    add_shape(slide, x, y, Inches(2.2), Pt(6), color)

    # 배지
    badge_shape = add_shape(slide, x + Inches(0.3), y + Inches(0.3), Inches(1.6), Inches(0.35), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(1.6), Inches(0.35),
                 badge, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 제품명
    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(1.8), Inches(0.9),
                 name, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 설명
    add_text_box(slide, x + Inches(0.2), y + Inches(2.0), Inches(1.8), Inches(1),
                 desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

# 하단 연결 화살표 느낌
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
             "모든 제품은 동일한 Claude AI 모델 (Opus 4.6 / Sonnet 4.6 / Haiku 4.5)을 기반으로 합니다",
             font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: Claude 생태계 관계도
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_cream_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.5), Inches(1.5))
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.7),
             "Claude 생태계 관계도", font_size=30, color=BLACK, bold=True)

# 중앙 Claude AI 모델
center_x, center_y = Inches(5.5), Inches(3.5)
center = add_circle(slide, center_x, center_y, Inches(1.8), CLAUDE_CORAL)
add_text_box(slide, center_x + Inches(0.15), center_y + Inches(0.5), Inches(1.5), Inches(0.8),
             "Claude AI\n모델", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 주변 제품들
satellites = [
    ("Claude.ai\n(Web/App)", Inches(2.5), Inches(1.2), CLAUDE_BLUE),
    ("Claude Code\n(CLI)", Inches(8.5), Inches(1.2), CLAUDE_GREEN),
    ("Claude\nDesktop", Inches(1.5), Inches(4.5), CLAUDE_PURPLE),
    ("Claude\nCowork", Inches(9.5), Inches(4.5), CLAUDE_ORANGE),
    ("OpenClaw\n(자가호스팅)", Inches(5.5), Inches(6), CLAUDE_CORAL),
]

for name, sx, sy, color in satellites:
    sat = add_circle(slide, sx, sy, Inches(1.3), color)
    add_text_box(slide, sx + Inches(0.1), sy + Inches(0.3), Inches(1.1), Inches(0.7),
                 name, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 연결 설명
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(2), Inches(0.4),
             "브라우저 접속", font_size=10, color=CLAUDE_BLUE, bold=True)
add_text_box(slide, Inches(10), Inches(1.3), Inches(2), Inches(0.4),
             "터미널 CLI", font_size=10, color=CLAUDE_GREEN, bold=True)
add_text_box(slide, Inches(0.2), Inches(5.9), Inches(2), Inches(0.4),
             "네이티브 앱", font_size=10, color=CLAUDE_PURPLE, bold=True)
add_text_box(slide, Inches(10.8), Inches(5.9), Inches(2), Inches(0.4),
             "업무 통합", font_size=10, color=CLAUDE_ORANGE, bold=True)
add_text_box(slide, Inches(7.2), Inches(6.8), Inches(2.5), Inches(0.4),
             "메신저 기반 24/7", font_size=10, color=CLAUDE_CORAL, bold=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: Claude AI (Web)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), CLAUDE_BLUE)
add_accent_line(slide, Inches(0.8), Inches(0.5), Inches(1.5), CLAUDE_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(5), Inches(0.7),
             "Claude AI (Web)", font_size=30, color=BLACK, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.4),
             "claude.ai  |  가장 쉽게 시작하는 Claude", font_size=14, color=CLAUDE_BLUE)

# 왼쪽: 특징
add_card(slide, Inches(0.8), Inches(2), Inches(3.8), Inches(4.5),
         "핵심 기능", [
             "  대화형 AI 어시스턴트",
             "  Projects로 컨텍스트 관리",
             "  파일 업로드 & 분석 (PDF, 이미지 등)",
             "  Artifacts (코드, 문서, 차트 생성)",
             "  Web Search (실시간 정보 검색)",
             "  Claude.ai/code (웹 기반 코딩)",
             "  MCP 서버 연결 (Notion, Figma 등)",
             "  모바일 앱 (iOS / Android)",
         ], icon_color=CLAUDE_BLUE)

# 중앙: 활용 팁
add_card(slide, Inches(5), Inches(2), Inches(3.8), Inches(4.5),
         "이런 분에게 추천", [
             "  AI를 처음 사용하는 입문자",
             "  빠른 질문 & 답변이 필요할 때",
             "  문서 요약 / 번역이 필요할 때",
             "  코드 없이 AI를 활용하고 싶을 때",
             "  팀 협업이 필요할 때 (Team 플랜)",
             "",
             "  Tip: Projects에 자주 쓰는 지시사항을",
             "  저장하면 매번 반복 입력 불필요!",
         ], icon_color=CLAUDE_BLUE)

# 오른쪽: 플랜
add_card(slide, Inches(9.2), Inches(2), Inches(3.5), Inches(4.5),
         "구독 플랜", [
             "  Free — 무료, 제한적 사용",
             "  Pro — $20/월, 5x 더 많은 사용량",
             "  Max — $100/월, 20x 사용량",
             "  Max+ — $200/월, 무제한에 가까움",
             "  Team — $25/인/월, 팀 협업",
             "  Enterprise — 맞춤형",
             "",
             "  Max 이상: Remote Control 사용 가능",
         ], icon_color=CLAUDE_BLUE)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: Claude Code (CLI)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), CLAUDE_GREEN)
add_shape(slide, Inches(0.8), Inches(0.5), Inches(1.5), Pt(4), CLAUDE_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(5), Inches(0.7),
             "Claude Code (CLI)", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(8), Inches(0.4),
             "터미널에서 직접 코딩하는 AI 에이전트  |  npm install -g @anthropic-ai/claude-code",
             font_size=13, color=CLAUDE_GREEN)

# 특징 카드들
features = [
    ("코드 읽기 & 수정", "프로젝트 전체를 이해하고\n코드를 직접 수정합니다", CLAUDE_GREEN),
    ("Git & GitHub", "커밋, PR 생성, 이슈 관리를\n자동으로 처리합니다", CLAUDE_BLUE),
    ("터미널 명령 실행", "빌드, 테스트, 배포 등\n쉘 명령을 직접 실행합니다", CLAUDE_PURPLE),
    ("MCP 서버 연동", "외부 도구와 연결하여\n확장된 기능을 사용합니다", CLAUDE_ORANGE),
]

for i, (title, desc, color) in enumerate(features):
    x = Inches(0.6) + Inches(i * 3.15)
    card = add_shape(slide, x, Inches(2), Inches(2.9), Inches(2), CLAUDE_NAVY)
    add_shape(slide, x, Inches(2), Inches(2.9), Pt(5), color)
    add_text_box(slide, x + Inches(0.2), Inches(2.3), Inches(2.5), Inches(0.5),
                 title, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.9), Inches(2.5), Inches(1),
                 desc, font_size=12, color=GRAY)

# 주요 명령어
cmd_box = add_shape(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(2.7), CLAUDE_NAVY)
add_text_box(slide, Inches(0.9), Inches(4.4), Inches(5), Inches(0.4),
             "자주 쓰는 명령어", font_size=16, color=CLAUDE_GREEN, bold=True)

cmds = [
    ("claude", "대화형 모드 시작"),
    ("/init", "프로젝트에 CLAUDE.md 생성"),
    ("/rc  또는  /remote-control", "Remote Control 활성화 (NEW!)"),
    ("/config", "설정 변경"),
    ("claude remote-control", "쉘에서 직접 원격 제어 시작"),
    ("claude --model opus", "모델 지정하여 실행"),
]

for i, (cmd, desc) in enumerate(cmds):
    col = Inches(0.9) if i < 3 else Inches(6.5)
    row = Inches(5) + Inches((i % 3) * 0.5)
    add_text_box(slide, col, row, Inches(2.5), Inches(0.4),
                 cmd, font_size=12, color=CLAUDE_GREEN, bold=True, font_name="Consolas")
    add_text_box(slide, col + Inches(2.8), row, Inches(3), Inches(0.4),
                 desc, font_size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: Claude Desktop & Cowork
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), CLAUDE_PURPLE)
add_accent_line(slide, Inches(0.8), Inches(0.5), Inches(1.5), CLAUDE_PURPLE)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.7),
             "Claude Desktop & Cowork", font_size=30, color=BLACK, bold=True)

# Desktop 섹션
add_card(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.5),
         "Claude Desktop (데스크톱 앱)", [
             "  macOS / Windows 네이티브 앱",
             "  MCP 서버 연결 가능 (파일시스템, GitHub 등)",
             "  바탕화면에서 바로 실행",
             "  Web 버전과 동일한 대화 기능",
             "  로컬 파일에 쉽게 접근",
         ], icon_color=CLAUDE_PURPLE)

# Cowork 섹션
add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.5),
         "Claude Cowork (업무 자동화)", [
             "  macOS 전용 (Desktop 내 기능)",
             "  실제 업무 도구를 AI가 직접 조작",
             "  Pro / Max 플랜 필요",
             "  11개 공식 플러그인 지원",
             "  '컴퓨터 사용' 기능 기반",
         ], icon_color=CLAUDE_ORANGE)

# Cowork 플러그인 목록
add_text_box(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(0.5),
             "Cowork 공식 플러그인 11종", font_size=16, color=CLAUDE_PURPLE, bold=True)

plugins = [
    ("Apple Mail", "이메일 읽기/작성"),
    ("Apple Notes", "메모 관리"),
    ("Apple Calendar", "일정 관리"),
    ("Apple Reminders", "할 일 관리"),
    ("Notion", "노트/DB 접근"),
    ("Google Docs", "문서 작성"),
    ("Google Sheets", "스프레드시트"),
    ("Google Slides", "프레젠테이션"),
    ("Excel / Word / PPT", "MS Office 연동"),
    ("Finder", "파일 탐색"),
    ("Safari", "웹 브라우징"),
]

for i, (name, desc) in enumerate(plugins):
    col = Inches(0.8) + Inches((i // 4) * 4)
    row = Inches(4.9) + Inches((i % 4) * 0.45)
    dot = add_circle(slide, col, row + Pt(4), Pt(8), CLAUDE_PURPLE)
    add_text_box(slide, col + Inches(0.2), row, Inches(1.6), Inches(0.4),
                 name, font_size=11, color=BLACK, bold=True)
    add_text_box(slide, col + Inches(1.8), row, Inches(1.5), Inches(0.4),
                 desc, font_size=11, color=GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: Remote Control (NEW!)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), CLAUDE_ORANGE)
add_shape(slide, Inches(0.8), Inches(0.5), Inches(1.5), Pt(4), CLAUDE_ORANGE)

# NEW 배지
badge = add_shape(slide, Inches(7), Inches(0.6), Inches(1.2), Inches(0.4), CLAUDE_ORANGE)
add_text_box(slide, Inches(7), Inches(0.6), Inches(1.2), Inches(0.4),
             "NEW!", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(0.7), Inches(6), Inches(0.7),
             "Remote Control", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "폰에서 PC의 Claude Code를 원격 조종  |  2026년 2월 공개",
             font_size=14, color=CLAUDE_ORANGE)

# 흐름도
flow_items = [
    ("PC에서 Claude Code 실행", "claude remote-control\n또는 /rc", CLAUDE_GREEN),
    ("QR 코드 / URL 생성", "스페이스바로 QR 토글", CLAUDE_BLUE),
    ("폰에서 스캔 & 연결", "Claude 앱 또는 브라우저", CLAUDE_ORANGE),
    ("어디서든 코딩!", "파일 수정, 테스트, Git 등", CLAUDE_CORAL),
]

for i, (title, desc, color) in enumerate(flow_items):
    x = Inches(0.5) + Inches(i * 3.2)
    card = add_shape(slide, x, Inches(2.2), Inches(2.8), Inches(1.8), CLAUDE_NAVY)
    add_shape(slide, x, Inches(2.2), Inches(2.8), Pt(5), color)

    # 번호
    num_circle = add_circle(slide, x + Inches(0.15), Inches(2.5), Inches(0.4), color)
    add_text_box(slide, x + Inches(0.15), Inches(2.5), Inches(0.4), Inches(0.4),
                 str(i+1), font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.6), Inches(2.5), Inches(2), Inches(0.5),
                 title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(3.2), Inches(2.4), Inches(0.6),
                 desc, font_size=11, color=GRAY)

    # 화살표 (마지막 제외)
    if i < 3:
        add_text_box(slide, x + Inches(2.85), Inches(2.8), Inches(0.3), Inches(0.5),
                     "→", font_size=20, color=CLAUDE_ORANGE, bold=True)

# 보안 & 요구사항
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.5),
             "보안 모델", font_size=16, color=CLAUDE_ORANGE, bold=True)

security_items = [
    "  인바운드 포트를 열지 않음 — 아웃바운드 HTTPS만 사용",
    "  모든 통신은 TLS 암호화 (Anthropic API 경유)",
    "  파일과 MCP 서버는 로컬에서만 처리",
    "  단기 수명 인증 토큰 사용",
]
lines = [{"text": item, "size": 12, "color": GRAY, "space_after": 2} for item in security_items]
add_multi_text(slide, Inches(0.8), Inches(5), Inches(5.5), Inches(2), lines)

# 요구사항
add_text_box(slide, Inches(7), Inches(4.5), Inches(5), Inches(0.5),
             "요구사항", font_size=16, color=CLAUDE_ORANGE, bold=True)

req_items = [
    "  Pro 또는 Max 플랜 필수",
    "  /login 으로 claude.ai 인증",
    "  Claude Code v2.1.52 이상",
    "  PC가 켜져있고 Claude Code 실행 중이어야 함",
]
lines = [{"text": item, "size": 12, "color": GRAY, "space_after": 2} for item in req_items]
add_multi_text(slide, Inches(7), Inches(5), Inches(5.5), Inches(2), lines)


# ══════════════════════════════════════════════════════════════
# SLIDE 9: 제품별 비교표
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_cream_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.5), Inches(1.5))
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.7),
             "제품별 비교표", font_size=30, color=BLACK, bold=True)

# 테이블 헤더
headers = ["기능", "Claude AI\n(Web)", "Claude\nCode", "Claude\nDesktop", "Claude\nCowork", "Open\nClaw"]
header_colors = [CLAUDE_DARK, CLAUDE_BLUE, CLAUDE_GREEN, CLAUDE_PURPLE, CLAUDE_ORANGE, CLAUDE_CORAL]

for i, (header, color) in enumerate(zip(headers, header_colors)):
    x = Inches(0.5) + Inches(i * 2.1)
    w = Inches(1.3) if i == 0 else Inches(1.9)
    if i == 0:
        w = Inches(2)
    h_shape = add_shape(slide, x, Inches(1.6), w, Inches(0.7), color)
    add_text_box(slide, x, Inches(1.6), w, Inches(0.7),
                 header, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 테이블 행
rows = [
    ("접근 방식", "브라우저/앱", "터미널 CLI", "데스크톱 앱", "데스크톱 앱", "메신저 봇"),
    ("대화 / 질답", "●", "●", "●", "●", "●"),
    ("코드 작성/수정", "△", "●", "△", "○", "△"),
    ("파일 직접 수정", "✕", "●", "MCP", "●", "●"),
    ("업무 도구 연동", "MCP", "MCP", "MCP", "직접 조작", "스킬/MCP"),
    ("24/7 상시 운영", "✕", "✕", "✕", "✕", "●"),
    ("모바일 접근", "●", "RC로 가능", "✕", "✕", "메신저"),
    ("무료 사용", "●", "✕", "●", "✕", "●(OSS)"),
    ("추천 대상", "입문자", "개발자", "일반 사용자", "업무 자동화", "파워유저"),
]

for r, row_data in enumerate(rows):
    y = Inches(2.35) + Inches(r * 0.52)
    bg_color = WHITE if r % 2 == 0 else CLAUDE_LIGHT
    for c, cell in enumerate(row_data):
        x = Inches(0.5) + Inches(c * 2.1)
        w = Inches(2) if c == 0 else Inches(1.9)
        cell_shape = add_shape(slide, x, y, w, Inches(0.5), bg_color)
        fc = BLACK if c == 0 else DARK_GRAY
        fb = True if c == 0 else False
        add_text_box(slide, x, y + Pt(4), w, Inches(0.4),
                     cell, font_size=10, color=fc, bold=fb, alignment=PP_ALIGN.CENTER)

# 범례
add_text_box(slide, Inches(0.8), Inches(7), Inches(10), Inches(0.3),
             "● = 완전 지원    ○ = 부분 지원    △ = 제한적    ✕ = 미지원    RC = Remote Control    MCP = 플러그인으로 확장 가능",
             font_size=10, color=GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 10: OpenClaw 소개
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), CLAUDE_CORAL)
add_shape(slide, Inches(0.8), Inches(0.5), Inches(1.5), Pt(4), CLAUDE_CORAL)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.7),
             "OpenClaw — 실행형 AI 에이전트", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             '"웹에서 답만 받는 AI"가 아니라 "내 환경에서 실제로 일하는 AI 비서"',
             font_size=14, color=CLAUDE_CORAL)

# 3개 핵심 카드
oc_features = [
    ("자가 호스팅", [
        "내 PC/서버에서 24/7 실행",
        "데이터가 외부로 나가지 않음",
        "WSL2 Ubuntu에서 운영",
        "무료 오픈소스 (MIT)",
    ], CLAUDE_CORAL),
    ("메신저 통합", [
        "WhatsApp / Telegram / Slack",
        "Discord / iMessage 지원",
        "카카오톡 (브릿지로 연동)",
        "어디서든 메시지로 AI 지시",
    ], CLAUDE_BLUE),
    ("실제 작업 수행", [
        "파일 정리 / 문서 요약",
        "브라우저 자동화",
        "크론 스케줄 자동화",
        "이메일 분류 / 알림",
    ], CLAUDE_GREEN),
]

for i, (title, items, color) in enumerate(oc_features):
    x = Inches(0.6) + Inches(i * 4.1)
    card = add_shape(slide, x, Inches(2.2), Inches(3.8), Inches(3), CLAUDE_NAVY)
    add_shape(slide, x, Inches(2.2), Inches(3.8), Pt(5), color)
    add_text_box(slide, x + Inches(0.3), Inches(2.5), Inches(3.2), Inches(0.5),
                 title, font_size=18, color=color, bold=True)
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.3), Inches(3.2) + Inches(j * 0.45), Inches(3.2), Inches(0.4),
                     f"  {item}", font_size=12, color=GRAY)

# 아키텍처
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(5), Inches(0.5),
             "기본 아키텍처", font_size=16, color=CLAUDE_CORAL, bold=True)

arch_text = "Gateway (포트 18789)  →  Agent (Claude Sonnet 4.6)  →  채널 (WhatsApp/Telegram)  →  사용자"
add_text_box(slide, Inches(0.8), Inches(6), Inches(11), Inches(0.4),
             arch_text, font_size=13, color=GRAY, font_name="Consolas")

# 빠른 시작
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
             "빠른 시작:  npm install -g openclaw@latest  →  npx openclaw onboard --install-daemon  →  npx openclaw tui",
             font_size=11, color=CLAUDE_GREEN, font_name="Consolas")


# ══════════════════════════════════════════════════════════════
# SLIDE 11: OpenClaw 카카오톡 연동
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), CLAUDE_ORANGE)
add_accent_line(slide, Inches(0.8), Inches(0.5), Inches(1.5), CLAUDE_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.7),
             "OpenClaw + 카카오톡 연동", font_size=30, color=BLACK, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(8), Inches(0.4),
             "github.com/tornado1014/openclaw-kakao  |  BlueStacks 기반 브릿지 시스템",
             font_size=13, color=CLAUDE_ORANGE)

# 아키텍처 다이어그램
arch_boxes = [
    ("BlueStacks 5\nKakaoTalk\n+ MessengerBotR", Inches(0.5), CLAUDE_ORANGE),
    ("Bridge Server\n(Node.js)\n포트: 8787", Inches(4.5), CLAUDE_BLUE),
    ("OpenClaw\nGateway\n(Claude AI)", Inches(8.5), CLAUDE_GREEN),
]

for text, x, color in arch_boxes:
    box = add_shape(slide, x, Inches(2), Inches(3.2), Inches(1.8), color)
    add_text_box(slide, x + Inches(0.2), Inches(2.2), Inches(2.8), Inches(1.4),
                 text, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 화살표
add_text_box(slide, Inches(3.8), Inches(2.6), Inches(0.6), Inches(0.5),
             "→", font_size=28, color=CLAUDE_CORAL, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.8), Inches(2.6), Inches(0.6), Inches(0.5),
             "→", font_size=28, color=CLAUDE_CORAL, bold=True, alignment=PP_ALIGN.CENTER)

# ADB Watcher
watcher = add_shape(slide, Inches(4.5), Inches(4.2), Inches(3.2), Inches(0.7), CLAUDE_PURPLE)
add_text_box(slide, Inches(4.5), Inches(4.2), Inches(3.2), Inches(0.7),
             "ADB Watcher (Python) — 이미지 자동 감지", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 주요 명령어 & 주의사항
add_card(slide, Inches(0.5), Inches(5.2), Inches(5.8), Inches(2),
         "사용 가능한 명령어", [
             "  /ping — 봇 상태 확인    /status — 시스템 상태",
             "  /clear — 세션 리셋       /help — 도움말",
             "  /질문 <메시지> — 그룹채팅에서 AI 질문",
             "  1:1 채팅은 바로 메시지 → AI 응답",
         ], icon_color=CLAUDE_ORANGE)

add_card(slide, Inches(6.7), Inches(5.2), Inches(6), Inches(2),
         "주의사항", [
             "  반드시 카카오톡 부계정 사용 (본계정 정지 위험)",
             "  BlueStacks ADB 디버깅 활성화 필수",
             "  이미지 분석: 카톡 '사진 자동 다운로드' 켜기",
             "  Node.js 18+ / Python 3.8+ 필요",
         ], icon_color=RGBColor(0xE0, 0x40, 0x40))


# ══════════════════════════════════════════════════════════════
# SLIDE 12: 자동화 활용 — 입문자
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_cream_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.5), Inches(1.5), CLAUDE_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.7),
             "자동화 활용 — 입문자 편", font_size=30, color=BLACK, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(8), Inches(0.4),
             "오늘 바로 시작할 수 있는 자동화 3가지", font_size=14, color=CLAUDE_GREEN)

# 3개 자동화
automations = [
    ("아침 브리핑", "매일 아침 일정 + 날씨 + 할 일 요약",
     'openclaw cron add \\\n  --name "morning-briefing" \\\n  --schedule "0 8 * * *" \\\n  --message "오늘 일정과 할 일 정리해줘" \\\n  --session isolated',
     CLAUDE_ORANGE),
    ("리마인더 봇", "중요한 일정 / 운동 / 약 알림",
     'openclaw cron add \\\n  --name "exercise-reminder" \\\n  --schedule "0 19 * * *" \\\n  --message "오늘 운동했나요? 30분만 걸어요!" \\\n  --session isolated',
     CLAUDE_BLUE),
    ("뉴스 요약", "관심 분야 뉴스를 매일 클리핑",
     'openclaw cron add \\\n  --name "news-digest" \\\n  --schedule "0 9 * * 1-5" \\\n  --message "AI/테크 최신 뉴스 3개 요약해줘" \\\n  --session isolated',
     CLAUDE_PURPLE),
]

for i, (title, desc, cmd, color) in enumerate(automations):
    x = Inches(0.5) + Inches(i * 4.2)
    # 제목
    num_c = add_circle(slide, x, Inches(1.9), Inches(0.5), color)
    add_text_box(slide, x, Inches(1.9), Inches(0.5), Inches(0.5),
                 str(i+1), font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.6), Inches(1.95), Inches(3), Inches(0.4),
                 title, font_size=18, color=BLACK, bold=True)
    add_text_box(slide, x + Inches(0.6), Inches(2.4), Inches(3.4), Inches(0.4),
                 desc, font_size=12, color=GRAY)
    # 코드 박스
    code_bg = add_shape(slide, x, Inches(2.9), Inches(3.9), Inches(2.5), CLAUDE_DARK)
    add_text_box(slide, x + Inches(0.2), Inches(3), Inches(3.5), Inches(2.3),
                 cmd, font_size=9, color=CLAUDE_GREEN, font_name="Consolas")

# 팁
add_shape(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2), CLAUDE_LIGHT)
add_text_box(slide, Inches(0.8), Inches(5.9), Inches(11), Inches(0.4),
             "Tip: --session isolated 를 항상 사용하세요!", font_size=14, color=CLAUDE_CORAL, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
             "세션을 격리하면 컨텍스트가 누적되지 않아 비용을 절약하고, 각 작업이 깨끗한 상태에서 시작됩니다.",
             font_size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 13: 자동화 활용 — 중급/고급
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_shape(slide, Inches(0.8), Inches(0.5), Inches(1.5), Pt(4), CLAUDE_PURPLE)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.7),
             "자동화 활용 — 중급/고급 편", font_size=30, color=WHITE, bold=True)

# 중급
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
             "중급 자동화", font_size=18, color=CLAUDE_BLUE, bold=True)

mid_items = [
    ("Webhook 연동", "GitHub PR → OpenClaw 코드 리뷰 → Slack 알림", CLAUDE_BLUE),
    ("멀티 에이전트", "역할별 에이전트 분리 (고객응대/마케팅/개인)", CLAUDE_GREEN),
    ("이메일 분류", "중요 메일 감지 → 알림 + 요약 전송", CLAUDE_PURPLE),
]

for i, (title, desc, color) in enumerate(mid_items):
    y = Inches(2) + Inches(i * 0.7)
    dot = add_circle(slide, Inches(0.8), y + Pt(4), Pt(10), color)
    add_text_box(slide, Inches(1.2), y, Inches(2.5), Inches(0.4),
                 title, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(3.8), y, Inches(8), Inches(0.4),
                 desc, font_size=12, color=GRAY)

# 고급
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5), Inches(0.5),
             "고급 자동화", font_size=18, color=CLAUDE_ORANGE, bold=True)

adv_items = [
    ("ClawDeck 대시보드", "모든 에이전트를 한 화면에서 모니터링/관리", CLAUDE_ORANGE),
    ("브라우저 자동화", "경쟁사 가격 모니터링, 데이터 수집", CLAUDE_CORAL),
    ("Ollama 로컬 모델", "민감 데이터는 로컬 AI로, 복잡한 작업은 Claude로", CLAUDE_GREEN),
    ("Custom Skill 개발", "나만의 스킬 만들어 ClawHub에 공유", CLAUDE_BLUE),
]

for i, (title, desc, color) in enumerate(adv_items):
    y = Inches(4.8) + Inches(i * 0.6)
    dot = add_circle(slide, Inches(0.8), y + Pt(4), Pt(10), color)
    add_text_box(slide, Inches(1.2), y, Inches(2.5), Inches(0.4),
                 title, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(3.8), y, Inches(8), Inches(0.4),
                 desc, font_size=12, color=GRAY)

# 모델 역할 분리
add_shape(slide, Inches(7.5), Inches(1.5), Inches(5.2), Inches(2.2), CLAUDE_NAVY)
add_text_box(slide, Inches(7.8), Inches(1.6), Inches(4.6), Inches(0.4),
             "모델 역할 분리 전략", font_size=14, color=CLAUDE_CORAL, bold=True)

model_roles = [
    ("Opus 4.6", "복잡한 분석/코딩", "$$$"),
    ("Sonnet 4.6", "일반 대화/응대", "$$"),
    ("Haiku 4.5", "간단 분류/요약", "$"),
    ("Ollama (로컬)", "개인정보 민감", "무료"),
]

for i, (model, role, cost) in enumerate(model_roles):
    y = Inches(2.1) + Inches(i * 0.38)
    add_text_box(slide, Inches(7.8), y, Inches(1.5), Inches(0.35),
                 model, font_size=10, color=CLAUDE_GREEN, bold=True, font_name="Consolas")
    add_text_box(slide, Inches(9.3), y, Inches(2), Inches(0.35),
                 role, font_size=10, color=GRAY)
    add_text_box(slide, Inches(11.5), y, Inches(1), Inches(0.35),
                 cost, font_size=10, color=CLAUDE_ORANGE, bold=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 14: 결합 활용 맵
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.5), Inches(1.5))
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.7),
             "결합 활용 맵 — 모든 것을 연결하기", font_size=30, color=BLACK, bold=True)

# 시나리오별 카드
scenarios = [
    ("프로젝트 개발", [
        "Claude Code로 코딩",
        "Remote Control로 외출 중 모니터링",
        "GitHub PR 자동 리뷰 (Webhook)",
        "Slack으로 빌드 결과 알림",
    ], CLAUDE_GREEN, "Claude Code\n+ Remote Control\n+ OpenClaw Webhook"),
    ("고객 응대", [
        "WhatsApp/카톡으로 고객 문의 수신",
        "AI가 1차 응답 (다국어 지원)",
        "복잡한 문의 → 이메일 정리",
        "ClawDeck에서 응대 현황 모니터링",
    ], CLAUDE_BLUE, "OpenClaw\n+ 카카오톡 브릿지\n+ ClawDeck"),
    ("마케팅 & SNS", [
        "주 3회 콘텐츠 아이디어 자동 생성",
        "해시태그 & 캡션 초안 작성",
        "트렌드 분석 (Brave Search)",
        "Notion에 콘텐츠 캘린더 관리",
    ], CLAUDE_PURPLE, "OpenClaw Cron\n+ Web Search\n+ MCP (Notion)"),
    ("학습 & 연구", [
        "논문/자료 자동 요약 (PDF)",
        "과제 마감 리마인더",
        "학습 계획 자동 수립",
        "질문 → Claude로 즉시 답변",
    ], CLAUDE_ORANGE, "Claude.ai\n+ OpenClaw Cron\n+ Telegram"),
]

for i, (title, items, color, tools) in enumerate(scenarios):
    x = Inches(0.4) + Inches(i * 3.2)

    # 카드
    card = add_shape(slide, x, Inches(1.6), Inches(3), Inches(4.5), WHITE)
    add_shape(slide, x, Inches(1.6), Inches(3), Pt(6), color)

    # 제목
    add_text_box(slide, x + Inches(0.2), Inches(1.85), Inches(2.6), Inches(0.5),
                 title, font_size=15, color=color, bold=True)

    # 항목
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.2), Inches(2.5) + Inches(j * 0.42),
                     Inches(2.6), Inches(0.4),
                     f"  {item}", font_size=10, color=DARK_GRAY)

    # 사용 도구
    tool_bg = add_shape(slide, x + Inches(0.15), Inches(4.5), Inches(2.7), Inches(1.2), CLAUDE_LIGHT)
    add_text_box(slide, x + Inches(0.25), Inches(4.55), Inches(2.5), Inches(0.3),
                 "사용 도구:", font_size=9, color=GRAY, bold=True)
    add_text_box(slide, x + Inches(0.25), Inches(4.85), Inches(2.5), Inches(0.8),
                 tools, font_size=9, color=color, bold=True)

# 하단 메시지
add_shape(slide, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.8), CLAUDE_DARK)
add_text_box(slide, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.6),
             "핵심: 하나의 도구가 아닌, 여러 도구를 조합하여 나만의 AI 워크플로를 구축하세요!",
             font_size=14, color=CLAUDE_CORAL, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 15: 프롬프트 팁
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_shape(slide, Inches(0.8), Inches(0.5), Inches(1.5), Pt(4), CLAUDE_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.7),
             "프롬프트 팁 & 베스트 프랙티스", font_size=30, color=WHITE, bold=True)

# DO / DON'T
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
             "DO — 이렇게 하세요", font_size=18, color=CLAUDE_GREEN, bold=True)

dos = [
    "구체적인 역할을 부여하세요: \"너는 고객 응대 전문가야\"",
    "출력 형식을 지정하세요: \"표 형태로 정리해줘\"",
    "단계별로 요청하세요: \"먼저 분석하고, 그 다음 제안해줘\"",
    "예시를 제공하세요: \"이런 형태로 만들어줘: [예시]\"",
    "CLAUDE.md / AGENTS.md에 지시사항을 저장하세요",
]

for i, tip in enumerate(dos):
    y = Inches(2.1) + Inches(i * 0.5)
    add_circle(slide, Inches(0.8), y + Pt(4), Pt(8), CLAUDE_GREEN)
    add_text_box(slide, Inches(1.2), y, Inches(5.5), Inches(0.45),
                 tip, font_size=11, color=GRAY)

add_text_box(slide, Inches(7), Inches(1.5), Inches(5), Inches(0.5),
             "DON'T — 피해야 할 것", font_size=18, color=RGBColor(0xE0, 0x50, 0x50), bold=True)

donts = [
    "너무 모호한 요청: \"뭔가 좋은 거 만들어줘\"",
    "한 번에 너무 많은 것을 요청",
    "중요한 작업을 검증 없이 자동 실행",
    "API 키나 비밀번호를 프롬프트에 포함",
    "크론을 너무 자주 설정 (비용 폭증)",
]

for i, tip in enumerate(donts):
    y = Inches(2.1) + Inches(i * 0.5)
    add_circle(slide, Inches(7), y + Pt(4), Pt(8), RGBColor(0xE0, 0x50, 0x50))
    add_text_box(slide, Inches(7.4), y, Inches(5.5), Inches(0.45),
                 tip, font_size=11, color=GRAY)

# OpenClaw 프롬프트 파일 구조
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.5),
             "OpenClaw Workspace 프롬프트 파일 구조", font_size=16, color=CLAUDE_ORANGE, bold=True)

ws_files = [
    ("AGENTS.md", "핵심 행동 규칙 & 안전 지침", "\"고객에게 가격 할인 약속 금지\""),
    ("SOUL.md", "성격 & 커뮤니케이션 스타일", "\"따뜻하고 전문적, 이모지 적절히\""),
    ("MEMORY.md", "장기 기억 (프로젝트, 선호도)", "\"Bella는 한/중/영 3개 국어 사용\""),
    ("HEARTBEAT.md", "자동 체크리스트 & 긴급도 분류", "\"20분 무응답 시 이메일 발송\""),
]

for i, (fname, desc, example) in enumerate(ws_files):
    y = Inches(5.4) + Inches(i * 0.45)
    add_text_box(slide, Inches(0.8), y, Inches(1.8), Inches(0.4),
                 fname, font_size=11, color=CLAUDE_GREEN, bold=True, font_name="Consolas")
    add_text_box(slide, Inches(2.8), y, Inches(3.5), Inches(0.4),
                 desc, font_size=11, color=GRAY)
    add_text_box(slide, Inches(6.5), y, Inches(6), Inches(0.4),
                 example, font_size=10, color=DARK_GRAY, font_name="맑은 고딕")


# ══════════════════════════════════════════════════════════════
# SLIDE 16: 비용 비교
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_cream_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.5), Inches(1.5), CLAUDE_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.7),
             "비용 비교 & 시작 가이드", font_size=30, color=BLACK, bold=True)

# 구독 비교
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
             "Claude 구독 플랜 비교", font_size=18, color=CLAUDE_CORAL, bold=True)

plans = [
    ("Free", "$0", "기본 대화, 제한적", GRAY),
    ("Pro", "$20/월", "5x 사용량, Remote Control", CLAUDE_BLUE),
    ("Max", "$100/월", "20x 사용량, Claude Code 포함", CLAUDE_GREEN),
    ("Max+", "$200/월", "무제한에 가까운 사용량", CLAUDE_PURPLE),
]

for i, (name, price, desc, color) in enumerate(plans):
    x = Inches(0.8) + Inches(i * 2.9)
    card = add_shape(slide, x, Inches(2), Inches(2.6), Inches(1.8), WHITE)
    add_shape(slide, x, Inches(2), Inches(2.6), Pt(5), color)
    add_text_box(slide, x + Inches(0.2), Inches(2.2), Inches(2.2), Inches(0.4),
                 name, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.65), Inches(2.2), Inches(0.4),
                 price, font_size=20, color=BLACK, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(3.15), Inches(2.2), Inches(0.5),
                 desc, font_size=10, color=GRAY)

# OpenClaw 비용
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5), Inches(0.5),
             "OpenClaw 운영 비용 (별도)", font_size=18, color=CLAUDE_CORAL, bold=True)

oc_costs = [
    ("OpenClaw 자체", "무료 (오픈소스)", CLAUDE_GREEN),
    ("라이트 사용", "$10~20/월 (API)", CLAUDE_BLUE),
    ("중간 사용", "$30~50/월 (API)", CLAUDE_ORANGE),
    ("헤비 사용", "$100+/월 (API)", RGBColor(0xE0, 0x50, 0x50)),
]

for i, (item, cost, color) in enumerate(oc_costs):
    y = Inches(4.8) + Inches(i * 0.45)
    dot = add_circle(slide, Inches(0.8), y + Pt(4), Pt(8), color)
    add_text_box(slide, Inches(1.2), y, Inches(2.5), Inches(0.4),
                 item, font_size=12, color=BLACK, bold=True)
    add_text_box(slide, Inches(3.8), y, Inches(3), Inches(0.4),
                 cost, font_size=12, color=color, bold=True)

# 비용 절약 팁
add_card(slide, Inches(7), Inches(4.2), Inches(5.5), Inches(2.8),
         "비용 절약 팁", [
             "  --session isolated 항상 사용 (컨텍스트 누적 방지)",
             "  간단한 작업은 Haiku 모델 사용 (10배 저렴)",
             "  크론 최소 간격 1시간 이상 유지",
             "  작업 완료 시 세션 리셋",
             "  민감 데이터는 Ollama 로컬 모델 (무료)",
         ], icon_color=CLAUDE_GREEN)


# ══════════════════════════════════════════════════════════════
# SLIDE 17: 시작하기 Step by Step
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_accent_line(slide, Inches(0.8), Inches(0.5), Inches(1.5), CLAUDE_CORAL)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.7),
             "오늘부터 시작하기", font_size=30, color=BLACK, bold=True)

steps = [
    ("1", "계정 만들기", "claude.ai 가입 → Pro 이상 구독", CLAUDE_BLUE,
     "5분"),
    ("2", "Claude Code 설치", "npm install -g @anthropic-ai/claude-code", CLAUDE_GREEN,
     "10분"),
    ("3", "Remote Control 설정", "/login → /rc → 폰으로 스캔", CLAUDE_ORANGE,
     "5분"),
    ("4", "OpenClaw 설치", "npm install -g openclaw@latest\nnpx openclaw onboard", CLAUDE_CORAL,
     "30분"),
    ("5", "첫 자동화 설정", "아침 브리핑 크론 등록", CLAUDE_PURPLE,
     "10분"),
    ("6", "확장하기", "멀티 에이전트, 카카오톡, Webhook...", CLAUDE_GREEN,
     "계속!"),
]

for i, (num, title, desc, color, time) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.2)
    y = Inches(1.5) + Inches(row * 2.8)

    # 카드
    card = add_shape(slide, x, y, Inches(3.9), Inches(2.3), WHITE)
    add_shape(slide, x, y, Inches(3.9), Pt(5), color)

    # 번호
    num_c = add_circle(slide, x + Inches(0.2), y + Inches(0.3), Inches(0.5), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(0.5), Inches(0.5),
                 num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 제목
    add_text_box(slide, x + Inches(0.8), y + Inches(0.35), Inches(2.5), Inches(0.4),
                 title, font_size=16, color=BLACK, bold=True)

    # 소요시간
    time_bg = add_shape(slide, x + Inches(2.8), y + Inches(0.35), Inches(0.9), Inches(0.35), color)
    add_text_box(slide, x + Inches(2.8), y + Inches(0.35), Inches(0.9), Inches(0.35),
                 time, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 설명
    add_text_box(slide, x + Inches(0.3), y + Inches(1), Inches(3.3), Inches(1),
                 desc, font_size=11, color=GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 18: 마무리
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

# 장식
add_circle(slide, Inches(-1), Inches(-1), Inches(4), CLAUDE_NAVY)
add_circle(slide, Inches(10), Inches(4.5), Inches(5), CLAUDE_NAVY)

add_shape(slide, Inches(4.5), Inches(1.8), Inches(0.8), Pt(5), CLAUDE_CORAL)

add_text_box(slide, Inches(2), Inches(2.2), Inches(9), Inches(1),
             "AI는 도구입니다.\n어떻게 조합하느냐가 실력입니다.",
             font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.8),
             "Claude AI + Claude Code + Remote Control + OpenClaw\n= 언제 어디서나 일하는 나만의 AI 팀",
             font_size=16, color=CLAUDE_CORAL, alignment=PP_ALIGN.CENTER)

# 참고 자료
add_shape(slide, Inches(2.5), Inches(5), Inches(8.3), Inches(1.5), CLAUDE_NAVY)
add_text_box(slide, Inches(2.8), Inches(5.1), Inches(7.7), Inches(0.3),
             "참고 자료", font_size=12, color=CLAUDE_CORAL, bold=True, alignment=PP_ALIGN.CENTER)

refs = [
    "claude.ai — Claude AI 웹",
    "code.claude.com/docs — Claude Code 공식 문서",
    "docs.openclaw.ai — OpenClaw 공식 문서",
    "github.com/tornado1014/openclaw-kakao — 카카오톡 연동",
]
for i, ref in enumerate(refs):
    add_text_box(slide, Inches(3), Inches(5.5) + Inches(i * 0.3), Inches(7), Inches(0.3),
                 ref, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(6.8), Inches(9), Inches(0.4),
             "Made with Claude Opus 4.6 by Bella (OZKIZ)  |  2026-02",
             font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════
output_path = os.path.join("D:/Claude_AI_Knowledge", "Claude_AI_Ecosystem_Guide_2026.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
